import streamlit as st
import fitz  # PyMuPDF for PDF extraction
from docx import Document
from pptx import Presentation
import io
import re

# Page configuration
st.set_page_config(
    page_title="SmartPrep AI",
    page_icon="",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better UI
st.markdown("""
<style>
    .stApp {
        background: linear-gradient(145deg, #0a0f1e 0%, #0c1222 100%);
    }
    .chat-message-user {
        background: linear-gradient(135deg, #3b82f6, #2563eb);
        color: white;
        padding: 12px 18px;
        border-radius: 18px;
        margin: 8px 0;
        max-width: 80%;
        align-self: flex-end;
        float: right;
        clear: both;
    }
    .chat-message-ai {
        background: rgba(30, 41, 70, 0.85);
        color: #f0f3ff;
        padding: 12px 18px;
        border-radius: 18px;
        margin: 8px 0;
        max-width: 80%;
        align-self: flex-start;
        border: 1px solid rgba(255,255,255,0.1);
        float: left;
        clear: both;
    }
    .stButton button {
    background: #3b82f6;
    color: white;
    border-radius: 25px;
    padding: 10px;
    font-weight: 500;
    transition: all 0.3s ease;

    display: flex;
    justify-content: center;
    align-items: center;
    text-align: center;
}
    .stButton button:hover {
        background: #2563eb;
        transform: translateY(-2px);
    }
    .stTextInput input {
        background: rgba(255,255,255,0.07);
        color: black;
        border-radius: 25px;
        border: 1px solid rgba(59,130,246,0.3);
    }
    .stTextInput input:focus {
        border-color: #3b82f6;
        box-shadow: 0 0 0 2px rgba(59,130,246,0.2);
    }
    .uploaded-file-info {
        background: rgba(59,130,246,0.2);
        padding: 10px;
        border-radius: 10px;
        margin: 10px 0;
        text-align: center;
        font-size: 0.9rem;
    }
    .mode-active {
        background: #3b82f6 !important;
        color: white !important;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
def init_session_state():
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {"role": "ai", "content": " **Hello! Welcome to SmartPrep AI!**\n\nI'm your intelligent study companion. Here's how to get started:\n\n📁 **Upload a document** (PDF, DOCX, PPTX, or TXT)\n🎯 **Choose a mode** (Exam Prep, Topic Prep, or Brush Up)\n💬 **Ask questions** about your study material\n\nLet's make learning smarter together! 🚀"}
        ]
    
    if "current_mode" not in st.session_state:
        st.session_state.current_mode = "exam"
    
    if "uploaded_file_content" not in st.session_state:
        st.session_state.uploaded_file_content = ""
    
    if "uploaded_file_name" not in st.session_state:
        st.session_state.uploaded_file_name = ""
    
    if "processing_complete" not in st.session_state:
        st.session_state.processing_complete = False

init_session_state()

# Function to extract text from different file types
def extract_text_from_file(uploaded_file):
    """Extract text from uploaded file"""
    file_extension = uploaded_file.name.split('.')[-1].lower()
    extracted_text = ""
    
    try:
        if file_extension == 'pdf':
            # Extract PDF text
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                extracted_text += page.get_text()
            pdf_document.close()

        elif file_extension == 'txt':
            extracted_text = uploaded_file.read().decode('utf-8')
        
        elif file_extension == 'docx':
            doc = Document(io.BytesIO(uploaded_file.read()))
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
        
        elif file_extension == 'pptx':
            prs = Presentation(io.BytesIO(uploaded_file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        else:
            return None, f"Unsupported file format: {file_extension}. Please upload PDF, DOCX, PPTX, or TXT files."
        
        if extracted_text and len(extracted_text.strip()) > 20:
            return extracted_text, None
        else:
            return None, "Could not extract text. The file might be image-based or empty."
            
    except Exception as e:
        return None, f"Error processing file: {str(e)}"

# Function to generate response based on mode and file content
def generate_response(user_question, mode, file_content, file_name):
    """Generate intelligent response based on user query"""
    
    if not file_content or len(file_content.strip()) < 50:
        return f" **No document loaded**\n\nPlease upload a document (PDF, DOCX, PPTX, or TXT) using the sidebar first. Once uploaded, I'll be able to answer your questions about the content."
    
    question_lower = user_question.lower()
    
    # Split into sentences
    sentences = re.split(r'[.!?]+', file_content)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    
    if not sentences:
        return " **No readable sentences found** in the document. The file might be image-based or poorly formatted."
    
    # Extract keywords from question (remove stop words)
    stop_words = {'what', 'is', 'are', 'the', 'a', 'an', 'of', 'to', 'in', 'for', 'on', 
                  'with', 'by', 'at', 'from', 'this', 'that', 'these', 'those', 'and', 
                  'or', 'but', 'so', 'because', 'if', 'then', 'else', 'when', 'where', 
                  'which', 'who', 'whom', 'whose', 'why', 'how', 'me', 'my', 'tell', 
                  'give', 'please', 'can', 'could', 'would', 'should'}
    
    keywords = [w for w in question_lower.split() if w not in stop_words and len(w) > 2]
    
    if not keywords:
        keywords = question_lower.split()[:3]
    
    # Find sentences containing keywords
    matching_sentences = []
    for sentence in sentences:
        sentence_lower = sentence.lower()
        match_count = 0
        matched_keywords = []
        for keyword in keywords:
            if keyword in sentence_lower:
                match_count += 1
                matched_keywords.append(keyword)
        if match_count > 0:
            matching_sentences.append({
                "sentence": sentence,
                "matches": match_count,
                "keywords_found": matched_keywords
            })
    
    # Sort by relevance
    matching_sentences.sort(key=lambda x: x["matches"], reverse=True)
    
    # Handle different question types
    if "summar" in question_lower or "summary" in question_lower or "overview" in question_lower:
        # Return summary (first 8 sentences)
        summary = ". ".join(sentences[:8])
        mode_context = ""
        if mode == "exam":
            mode_context = "\n\n **Exam Tip:** Focus on understanding these key points as they're likely to appear in tests."
        elif mode == "topic":
            mode_context = "\n\n **Topic Focus:** This summary covers the main concepts. Ask me to elaborate on any specific point."
        elif mode == "brushup":
            mode_context = "\n\n **Quick Review:** This is your rapid revision summary for quick recall."
        
        return f" **Summary of '{file_name}':**\n\n{summary}...{mode_context}"
    
    elif "key point" in question_lower or "important" in question_lower or "main" in question_lower or "essential" in question_lower:
        # Return important sentences
        important_keywords = ['important', 'key', 'main', 'significant', 'essential', 'critical', 'major', 'primary', 'fundamental']
        important_sentences = []
        for sentence in sentences:
            sent_lower = sentence.lower()
            if any(kw in sent_lower for kw in important_keywords):
                important_sentences.append(sentence)
            if len(important_sentences) >= 8:
                break
        
        if not important_sentences:
            important_sentences = sentences[:6]
        
        result = f" **Key Points from '{file_name}':**\n\n"
        for i, sent in enumerate(important_sentences[:8], 1):
            result += f"{i}. {sent}.\n\n"
        
        if mode == "exam":
            result += "\n **Exam Tip:** These key points are critical for your preparation. Make sure you understand each one."
        elif mode == "brushup":
            result += "\n **Quick Recall:** Review these points regularly for better retention."
        
        return result
    
    elif "list" in question_lower or "topics" in question_lower or "chapters" in question_lower or "contents" in question_lower:
        # Extract potential headings or topics
        headings = re.findall(r'^[A-Z][A-Z\s]{3,}|^[0-9]+[\.\)]\s+[A-Z]|^[A-Z][a-z]+[\s]+[A-Z]', file_content, re.MULTILINE)
        if headings:
            result = f" **Topics found in '{file_name}':**\n\n"
            for i, h in enumerate(headings[:15], 1):
                result += f"{i}. {h.strip()}\n"
            return result
        else:
            # Return first sentences as topics
            result = f" **Content preview from '{file_name}':**\n\n"
            for i, s in enumerate(sentences[:10], 1):
                preview = s[:80] + "..." if len(s) > 80 else s
                result += f"{i}. {preview}\n\n"
            return result
    
    elif "what" in question_lower or "explain" in question_lower or "describe" in question_lower or "tell me" in question_lower:
        # Return matching content
        if matching_sentences:
            result = f" **Answer from '{file_name}':**\n\n"
            for i, match in enumerate(matching_sentences[:5], 1):
                result += f"{i}. {match['sentence']}.\n"
                if match['keywords_found']:
                    result += f"   *(Matched: {', '.join(match['keywords_found'])})*\n\n"
            return result
        else:
            # No matches found - show preview
            preview = file_content[:600]
            return f" **Question:** \"{user_question}\"\n\n**No exact match found in the document.**\n\n**Here's what your document contains:**\n\n\"{preview}...\"\n\n**Try asking with different keywords or:**'\n• 'List the main topics'"
    
    elif "how many" in question_lower or "count" in question_lower:
        # Count pages or sections
        word_count = len(file_content.split())
        char_count = len(file_content)
        sentence_count = len(sentences)
        
        return f" **Document Statistics for '{file_name}':**\n\n• **Words:** {word_count:,}\n• **Characters:** {char_count:,}\n• **Sentences:** {sentence_count}\n• **Mode:** {mode.upper()}\n\nWhat would you like to know more about?"
    
    else:
        # Default: Return matching sentences or document preview
        if matching_sentences:
            result = f" **From '{file_name}':**\n\n"
            for i, match in enumerate(matching_sentences[:4], 1):
                result += f"{i}. {match['sentence']}.\n\n"
            
            # Add mode-specific guidance
            if mode == "exam":
                result += "\n **Exam Mode:** Need practice questions? Ask me to 'generate questions' for this topic."
            elif mode == "topic":
                result += "\n **Topic Mode:** Want deeper explanation? Ask me to 'explain' any concept in detail."
            elif mode == "brushup":
                result += "\n **Brush Up Mode:** For quick revision, ask for 'key points' or 'summarize"
            
            return result
        else:
            # Show document preview
            preview = file_content[:500]
            return f" **Document Preview:**\n\n\"{preview}...\"\n\n**No direct match for:** \"{user_question}\"\n\n**Try asking:**\n• 'What is this document about?'\n• 'Summarize this document'\n• 'What are the key points?'"

# Sidebar
with st.sidebar:
    st.markdown("#  SmartPrep AI")
    st.markdown("---")
    
    st.markdown("##  Modes")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button(" Exam Prep", use_container_width=True):
            st.session_state.current_mode = "exam"
            st.session_state.messages.append({"role": "ai", "content": " **Exam Prep Mode Activated!**\n\nI'll help you prepare for tests with:\n• Practice questions\n• Key concepts\n• Important points to remember\n\nUpload a document and ask me to generate questions or explain concepts!"})
            st.rerun()
    
    with col2:
        if st.button(" Topic Prep", use_container_width=True):
            st.session_state.current_mode = "topic"
            st.session_state.messages.append({"role": "ai", "content": " **Topic Prep Mode Activated!**\n\nI'll provide detailed explanations and comprehensive coverage of your study material.\n\nUpload a document and ask me to explain any topic in detail!"})
            st.rerun()
    
    with col3:
        if st.button(" Brush Up", use_container_width=True):
            st.session_state.current_mode = "brushup"
            st.session_state.messages.append({"role": "ai", "content": " **Brush Up Mode Activated!**\n\nI'll give you:\n• Quick summaries\n• Key points for revision\n• Rapid review materials\n\nPerfect for last-minute preparation!"})
            st.rerun()
    
    st.markdown("---")
    
    st.markdown("##  File Upload")
    
    # File uploader
    uploaded_file = st.file_uploader(
        "Upload your study material",
        type=['pdf', 'txt', 'docx', 'pptx'],
        help="Supported formats: PDF, TXT, DOCX, PPTX"
    )
    
    if uploaded_file is not None:
        if st.button(" Process Document", use_container_width=True, type="primary"):
            with st.spinner(" Processing document..."):
                extracted_text, error = extract_text_from_file(uploaded_file)
                
                if error:
                    st.error(error)
                    st.session_state.messages.append({"role": "ai", "content": f" {error}"})
                else:
                    st.session_state.uploaded_file_content = extracted_text
                    st.session_state.uploaded_file_name = uploaded_file.name
                    st.session_state.processing_complete = True
                    
                    preview = extracted_text[:300] if len(extracted_text) > 300 else extracted_text
                    
                    success_msg = f""" **Successfully loaded '{uploaded_file.name}'**

 **Document Statistics:**
• Characters: {len(extracted_text):,}
• Words: {len(extracted_text.split()):,}
• Sentences: {len(re.split(r'[.!?]+', extracted_text))}

**Preview:**
"{preview}..."

**You can now ask me questions about this document in {st.session_state.current_mode.upper()} mode!"""
                    
                    st.session_state.messages.append({"role": "ai", "content": success_msg})
                    st.rerun()
    
    # Show current file status
    if st.session_state.uploaded_file_name:
        st.markdown("---")
        st.markdown("###  Active Document")
        st.info(f"**{st.session_state.uploaded_file_name[:30]}**\n\n{len(st.session_state.uploaded_file_content):,} characters")
        
        if st.button(" Clear Document", use_container_width=True):
            st.session_state.uploaded_file_content = ""
            st.session_state.uploaded_file_name = ""
            st.session_state.processing_complete = False
            st.session_state.messages.append({"role": "ai", "content": " Document cleared. You can upload a new document anytime."})
            st.rerun()
    
    st.markdown("---")
    
    # Current mode display
    mode_icons = {"exam": "", "topic": "", "brushup": ""}
    mode_names = {"exam": "Exam Prep", "topic": "Topic Prep", "brushup": "Brush Up"}
    st.markdown(f"###  Current Mode")
    st.markdown(f"**{mode_icons[st.session_state.current_mode]} {mode_names[st.session_state.current_mode]}**")
    
    st.markdown("---")
    
    if st.button(" Clear Chat History", use_container_width=True):
        st.session_state.messages = [
            {"role": "ai", "content": " **Chat history cleared!**\n\nReady to continue studying. Upload a document and ask me anything!"}
        ]
        st.rerun()

# Main chat area
st.markdown("##  SmartPrep AI Assistant")

# Create a container for chat messages
chat_container = st.container()

with chat_container:
    # Display chat messages
    for message in st.session_state.messages:
        if message["role"] == "user":
            st.markdown(f'<div class="chat-message-user"> {message["content"]}</div>', unsafe_allow_html=True)
        else:
            st.markdown(f'<div class="chat-message-ai"> {message["content"]}</div>', unsafe_allow_html=True)

# Chat input area
st.markdown("---")

# Create a form for chat input
with st.form(key="chat_form", clear_on_submit=True):
    col1, col2 = st.columns([5, 1])
    
    with col1:
        user_input = st.text_input(
            "Ask a question about your document...",
            key="user_input",
            label_visibility="collapsed",
            placeholder=" Type your question here... (e.g., 'Summarize this document', 'What are the key points?', 'Explain the main concept')"
        )
    
    with col2:
        submit_button = st.form_submit_button(" Send", use_container_width=True)
    
    if submit_button and user_input:
        # Add user message
        st.session_state.messages.append({"role": "user", "content": user_input})
        
        # Generate AI response
        with st.spinner(" Analyzing document..."):
            response = generate_response(
                user_input,
                st.session_state.current_mode,
                st.session_state.uploaded_file_content,
                st.session_state.uploaded_file_name
            )
        
        # Add AI response
        st.session_state.messages.append({"role": "ai", "content": response})
        st.rerun()

# Footer
st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #64748b; font-size: 0.8rem;'>"
    " SmartPrep AI - Your Intelligent Study Companion | Powered by AI"
    "</div>",
    unsafe_allow_html=True
)